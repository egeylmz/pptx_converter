from pptx import Presentation
from typing import List, Dict


def extract_text_from_pptx(file_path: str) -> List[Dict]:
    """
    Extracts text from each slide in a PowerPoint file.

    Args:
        file_path: Path to the PowerPoint file

    Returns:
        List containing a dict for each slide. Each dict:
        {
            'slide_number': int,
            'text': str (all text combined),
            'text_blocks': List[str] (text paragraph by paragraph)
        }
    """
    try:
        prs = Presentation(file_path)
        slides_data = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            text_blocks = []
            all_text = []

            # Check each shape (textbox, placeholder, etc.)
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content = shape.text.strip()
                    text_blocks.append(text_content)
                    all_text.append(text_content)

            # If text exists, add to slide data
            if all_text:
                slides_data.append({
                    'slide_number': slide_num,
                    'text': '\n'.join(all_text),
                    'text_blocks': text_blocks
                })
            else:
                # If no text, add as empty slide
                slides_data.append({
                    'slide_number': slide_num,
                    'text': '',
                    'text_blocks': []
                })

        return slides_data

    except Exception as e:
        raise Exception(f"Error occurred while reading PowerPoint file: {str(e)}")


def get_slide_count(file_path: str) -> int:
    """
    Returns the total number of slides in the PowerPoint file.
    """
    try:
        prs = Presentation(file_path)
        return len(prs.slides)
    except Exception as e:
        raise Exception(f"Error occurred while reading PowerPoint file: {str(e)}")